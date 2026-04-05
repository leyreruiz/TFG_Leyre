"""Script to ingest all test topic files into ChromaDB.

Usage:
  python -m backend.ingest_topics           # ingest (keeps existing data)
  python -m backend.ingest_topics --clean   # wipe collection first, then ingest
"""

import logging
import os
import sys

try:
    import fitz
except Exception:
    fitz = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

sys.path.insert(0, os.path.abspath(os.path.join(os.path.dirname(__file__), "..")))

from backend.clients.bbdd_client import save_text_chroma, delete_collection

logger = logging.getLogger(__name__)


def _divide_in_chunks(text: str, chunk_size: int = 500, overlap: int = 50) -> list[str]:
    """Split a text into overlapping chunks."""
    chunks = []
    start = 0
    while start < len(text):
        end = min(start + chunk_size, len(text))
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        if end == len(text):
            break
        start = end - overlap
    return chunks


def _infer_topic_from_path_or_metadata(file_path: str, metadata: dict | None = None) -> str:
    if isinstance(metadata, dict):
        meta_topic = metadata.get("topic")
        if isinstance(meta_topic, str) and meta_topic.strip():
            return meta_topic.strip()
    return os.path.splitext(os.path.basename(file_path))[0]


def ingest_file_txt(file_path: str, metadata: dict | None = None) -> list[str]:
    """Read a .txt file and ingest it into ChromaDB in chunks."""
    if not os.path.exists(file_path):
        logger.error("File not found: %s", file_path)
        return []

    try:
        with open(file_path, "r", encoding="utf-8") as f:
            text = f.read()
    except Exception as e:
        logger.error("Error reading file: %s", e)
        return []

    logger.info("Reading %s (%d characters)...", file_path, len(text))
    chunks = _divide_in_chunks(text)
    logger.info("Split into %d chunks", len(chunks))

    ids_saved = []
    topic_name = _infer_topic_from_path_or_metadata(file_path, metadata)
    filename = os.path.basename(file_path)

    for i, chunk in enumerate(chunks):
        meta = dict(metadata or {})
        meta.update({"chunk": i, "source": filename})
        doc_id = save_text_chroma(chunk, metadata=meta, topic=topic_name)
        if doc_id:
            ids_saved.append(doc_id)

    logger.info("Ingest complete: %d chunks saved", len(ids_saved))
    return ids_saved


def ingest_file_pdf(file_path: str, metadata: dict | None = None) -> list[str]:
    """Read a .pdf file and ingest it into ChromaDB in chunks."""
    if fitz is None:
        logger.error("PyMuPDF is not installed. Install 'pymupdf' to ingest PDFs.")
        return []

    if not os.path.exists(file_path):
        logger.error("File not found: %s", file_path)
        return []

    try:
        with fitz.open(file_path) as pdf:
            pages = [page.get_text("text") for page in pdf]
        text = "\n".join(pages)
    except Exception as e:
        logger.error("Error reading PDF file: %s", e)
        return []

    if not text.strip():
        logger.error("PDF has no extractable text: %s", file_path)
        return []

    logger.info("Reading %s (%d characters extracted)...", file_path, len(text))
    chunks = _divide_in_chunks(text)
    logger.info("Split into %d chunks", len(chunks))

    ids_saved = []
    topic_name = _infer_topic_from_path_or_metadata(file_path, metadata)
    filename = os.path.basename(file_path)

    for i, chunk in enumerate(chunks):
        meta = dict(metadata or {})
        meta.update({"chunk": i, "source": filename})
        doc_id = save_text_chroma(chunk, metadata=meta, topic=topic_name)
        if doc_id:
            ids_saved.append(doc_id)

    logger.info("Ingest complete: %d chunks saved", len(ids_saved))
    return ids_saved


def ingest_file_pptx(file_path: str, metadata: dict | None = None) -> list[str]:
    """Read a .pptx file and ingest it into ChromaDB in chunks."""
    if Presentation is None:
        logger.error("python-pptx is not installed. Install 'python-pptx' to ingest PowerPoints.")
        return []

    if not os.path.exists(file_path):
        logger.error("File not found: %s", file_path)
        return []

    try:
        presentation = Presentation(file_path)
        text_parts = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_parts.append(shape.text)
                if hasattr(shape, "table") and shape.table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text:
                                text_parts.append(cell.text)
        text = "\n".join(text_parts)
    except Exception as e:
        logger.error("Error reading PowerPoint file: %s", e)
        return []

    if not text.strip():
        logger.error("PowerPoint has no extractable text: %s", file_path)
        return []

    logger.info("Reading %s (%d characters extracted)...", file_path, len(text))
    chunks = _divide_in_chunks(text)
    logger.info("Split into %d chunks", len(chunks))

    ids_saved = []
    topic_name = _infer_topic_from_path_or_metadata(file_path, metadata)
    filename = os.path.basename(file_path)

    for i, chunk in enumerate(chunks):
        meta = dict(metadata or {})
        meta.update({"chunk": i, "source": filename})
        doc_id = save_text_chroma(chunk, metadata=meta, topic=topic_name)
        if doc_id:
            ids_saved.append(doc_id)

    logger.info("Ingest complete: %d chunks saved", len(ids_saved))
    return ids_saved


def ingest_file(file_path: str, metadata: dict | None = None) -> list[str]:
    """Ingest a supported file (.txt, .pdf or .pptx) into ChromaDB."""
    extension = os.path.splitext(file_path)[1].lower()

    if extension == ".txt":
        return ingest_file_txt(file_path, metadata=metadata)
    if extension == ".pdf":
        return ingest_file_pdf(file_path, metadata=metadata)
    if extension == ".pptx":
        return ingest_file_pptx(file_path, metadata=metadata)

    logger.error("Unsupported file type '%s'. Use .txt, .pdf or .pptx", extension)
    return []


DATA_DIR = os.path.join(os.path.dirname(__file__), "data")

TOPICS = [
    ("neural_networks.txt", {"topic": "neural_networks", "course": "IA"}),
    ("data_bases.txt", {"topic": "data_bases", "course": "BBDD"}),
]


def main():
    logging.basicConfig(level=logging.INFO)
    logger.info("=== Topic ingest ===")

    for filename, metadata in TOPICS:
        path = os.path.join(DATA_DIR, filename)
        if not os.path.exists(path):
            logger.warning("Not found: %s", path)
            continue

        logger.info("--- Ingesting: %s ---", filename)
        ids = ingest_file(path, metadata=metadata)
        logger.info("  -> %d chunks saved.", len(ids))

    logger.info("=== Ingest complete ===")


if __name__ == "__main__":
    if "--clean" in sys.argv:
        logger.info("Wiping existing collection...")
        delete_collection()
    main()
